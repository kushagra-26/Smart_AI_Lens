import streamlit as st
import cv2
from PIL import Image
import imagehash
import pytesseract
import os
import tempfile
import yt_dlp as youtube_dl
from pptx import Presentation
from pptx.util import Inches

# Configure Tesseract path if on Windows
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Function to download the video from YouTube
def download_video(url, output_path='video.mp4'):
    try:
        ydl_opts = {
            'format': 'best',
            'outtmpl': output_path
        }
        with youtube_dl.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])
        return output_path
    except Exception as e:
        st.error("Failed to download video. Please check the URL or try again later.")
        print(f"Error: {e}")
        return None

# Function to extract frames from the video at specific intervals
def extract_frames(video_path, frame_interval=30):
    cap = cv2.VideoCapture(video_path)
    frame_id = 0
    frames = []
    
    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break
        if frame_id % frame_interval == 0:
            frames.append(frame)
        frame_id += 1
    cap.release()
    return frames

# Function to get unique slides based on perceptual hashing
def get_unique_slides(frames, hash_difference_threshold=5):
    hashes = []
    unique_frames = []
    
    for frame in frames:
        pil_image = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
        hash_val = imagehash.phash(pil_image)
        
        # Only add frames that are sufficiently different from previous frames
        if all(abs(hash_val - h) > hash_difference_threshold for h in hashes):
            unique_frames.append(frame)
            hashes.append(hash_val)
    
    return unique_frames

# Function to extract text from each unique slide
def extract_text_from_slides(slides):
    slide_texts = []
    
    for idx, slide in enumerate(slides):
        pil_image = Image.fromarray(cv2.cvtColor(slide, cv2.COLOR_BGR2RGB))
        text = pytesseract.image_to_string(pil_image)
        
        # Convert to bytes for Streamlit compatibility
        slide_image = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        pil_image.save(slide_image.name)
        
        slide_texts.append((slide_image.name, text.strip()))
    
    return slide_texts

# Function to create a PowerPoint presentation from slides
def create_ppt(slide_texts):
    prs = Presentation()
    
    for slide_filename, text in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout
        left = top = Inches(1)
        slide.shapes.add_picture(slide_filename, left, top, width=Inches(8.5), height=Inches(6))
        txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8.5), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = text
    
    # Save PPT to a temporary file
    ppt_temp = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
    prs.save(ppt_temp.name)
    return ppt_temp.name

# Streamlit Interface
def main():
    st.title("YouTube Video Slide Extractor")
    st.write("Enter a YouTube video link, and this app will extract unique slides along with any text on them.")

    youtube_url = st.text_input("YouTube Video URL:", "")
    
    if youtube_url:
        if st.button("Process Video"):
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_video_path = os.path.join(temp_dir, "video.mp4")
                
                with st.spinner("Downloading video..."):
                    video_path = download_video(youtube_url, output_path=temp_video_path)
                    if not video_path:
                        st.error("Video download failed. Please check the URL.")
                        return

                with st.spinner("Extracting frames..."):
                    frames = extract_frames(video_path, frame_interval=30)

                with st.spinner("Detecting unique slides..."):
                    unique_slides = get_unique_slides(frames)
                
                with st.spinner("Extracting text from slides..."):
                    slide_texts = extract_text_from_slides(unique_slides)

                ppt_filename = create_ppt(slide_texts)

                st.success("Processing complete!")
                st.write("Extracted Slides and Text:")
                
                for slide_filename, text in slide_texts:
                    st.image(slide_filename, caption="Slide Image", use_column_width=True)
                    st.write("Extracted Text:")
                    st.text(text)
                    st.write("---")
                
                # Download link for the PowerPoint file
                with open(ppt_filename, "rb") as f:
                    st.download_button("Download PowerPoint Presentation", f, file_name="presentation.pptx")

if __name__ == "__main__":
    main()
